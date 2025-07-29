import argparse
import json
import os
import sys
import subprocess
import importlib.util
import tempfile
import re
import math
import concurrent.futures

def check_and_install_dependencies():
    """Check and install required dependencies."""
    required_packages = {
        'google-generativeai': 'google.generativeai',
        'pypdf2': 'PyPDF2',
        'python-docx': 'docx',
        'python-pptx': 'pptx'
    }
    
    missing_packages = []
    
    for package_name, import_name in required_packages.items():
        if importlib.util.find_spec(import_name) is None:
            missing_packages.append(package_name)
    
    if missing_packages:
        print(f"Installing missing dependencies: {', '.join(missing_packages)}")
        try:
            cmd = [sys.executable, "-m", "pip", "install", "--index-url", "https://pypi.org/simple/"] + missing_packages
            subprocess.check_call(cmd)
            print("Dependencies installed successfully")
        except subprocess.CalledProcessError as e:
            print(f"Error installing dependencies: {e}")
            print("\nTrying alternative installation method...")
            
            for package in missing_packages:
                try:
                    subprocess.check_call([sys.executable, "-m", "pip", "install", package])
                    print(f"Installed {package}")
                except subprocess.CalledProcessError:
                    print(f"Failed to install {package}. You may need to install it manually.")
                    if package == 'pypdf2':
                        print("Try: pip install PyPDF2==3.0.1")
                    elif package == 'python-docx':
                        print("Try: pip install python-docx==0.8.11")
                    elif package == 'python-pptx':
                        print("Try: pip install python-pptx==0.6.21")

    try:
        global genai, PyPDF2, Document, Presentation
        import google.generativeai as genai
        import PyPDF2
        from docx import Document
        from pptx import Presentation
        return True
    except ImportError as e:
        print(f"Failed to import required packages: {e}")
        print("Please install the required packages manually and try again.")
        return False

def configure_api(api_key):
    """Configure the Gemini API with the provided key."""
    try:
        genai.configure(api_key=api_key)
        return genai.GenerativeModel('gemini-1.5-flash-latest')  
    except Exception as e:
        print(f"Error configuring API: {e}")
        return None

def extract_text_from_pdf(pdf_path, page_range=None):
    """Extract text from a PDF file with optional page range."""
    text = ""
    try:
        with open(pdf_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            
            if page_range:
                start, end = page_range
                # Adjust for zero-based indexing
                start = max(0, start - 1)
                end = min(len(reader.pages), end)
                pages_to_read = range(start, end)
            else:
                pages_to_read = range(len(reader.pages))
            
            for page_num in pages_to_read:
                page_text = reader.pages[page_num].extract_text()
                if page_text:
                    text += page_text + "\n\n"
        
        if not text.strip():
            print("Warning: Extracted text is empty. The PDF might be scanned or image-based.")
        
        return text
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
        print("If using PyPDF2 version > 3.0.0, make sure you're using the updated API.")
        return None

def extract_text_from_txt(txt_path):
    """Extract text from a plain text file."""
    try:
        with open(txt_path, 'r', encoding='utf-8') as file:
            return file.read()
    except UnicodeDecodeError:
        try:
            with open(txt_path, 'r', encoding='latin-1') as file:
                return file.read()
        except Exception as e:
            print(f"Error reading text file with latin-1 encoding: {e}")
            return None
    except Exception as e:
        print(f"Error reading text file: {e}")
        return None

def extract_text_from_docx(docx_path):
    """Extract text from a Word document."""
    text = ""
    try:
        doc = Document(docx_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text
    except Exception as e:
        print(f"Error extracting text from DOCX: {e}")
        return None

def extract_text_from_pptx(pptx_path):
    """Extract text from a PowerPoint presentation."""
    text = ""
    try:
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n\n"
        return text
    except Exception as e:
        print(f"Error extracting text from PPTX: {e}")
        return None

def extract_chapters(text, chapter_names=None):
    """
    Extract specified chapters from the text.
    If chapter_names is None, return the full text.
    """
    if not chapter_names:
        return text
    
    chapter_patterns = [
        r'(?i)chapter\s+\d+[\s:]*(.+)', 
        r'(?i)unit\s+\d+[\s:]*(.+)',     
        r'(?i)section\s+\d+[\.\d]*[\s:]*(.+)', 
        r'(?i)module\s+\d+[\s:]*(.+)'    
    ]
    
    chapters = {}
    current_chapter = None
    current_content = []
    
    lines = text.split('\n')
    for line in lines:
        
        is_chapter_heading = False
        
        if chapter_names and any(chapter.lower() in line.lower() for chapter in chapter_names):
            current_chapter = line.strip()
            current_content = []
            is_chapter_heading = True
            continue
        
        for pattern in chapter_patterns:
            if re.search(pattern, line.strip()):
                if current_chapter and current_content:
                    chapters[current_chapter] = '\n'.join(current_content)
                
                current_chapter = line.strip()
                current_content = []
                is_chapter_heading = True
                break
        
        if not is_chapter_heading and current_chapter:
            current_content.append(line)
    
    if current_chapter and current_content:
        chapters[current_chapter] = '\n'.join(current_content)
    
    if chapter_names:
        filtered_chapters = {}
        for name in chapter_names:
            for chapter_title, content in chapters.items():
                if name.lower() in chapter_title.lower():
                    filtered_chapters[chapter_title] = content
        
        if filtered_chapters:
            return '\n\n'.join([f"{title}\n{content}" for title, content in filtered_chapters.items()])
    
    if chapters:
        return '\n\n'.join([f"{title}\n{content}" for title, content in chapters.items()])
    
    print("Warning: No chapters were detected in the text.")
    return text

def generate_question_bank(model, content, difficulty="medium", format_json=True):
    """Generate questions from the textbook content using Gemini API with CO numbers and Bloom's taxonomy levels."""
    difficulty_descriptions = {
        "easy": "suitable for beginners with basic understanding",
        "medium": "requiring moderate understanding of concepts",
        "hard": "challenging questions testing deep understanding",
        "mixed": "a combination of easy, medium, and hard questions"
    }
    
    difficulty_desc = difficulty_descriptions.get(difficulty.lower(), "requiring moderate understanding of concepts")
    
    max_content_length = 25000  
    if len(content) > max_content_length:
        print(f"Warning: Content is too long ({len(content)} chars). Truncating to {max_content_length} chars.")
        content = content[:max_content_length] + "\n\n[Content truncated due to length...]"
    
    course_outcomes = """
    Course Outcomes: After completing the course, the students will be able to:
    CO1: Elucidate the principles of management theory & recognize the characteristics of an organization.
    CO2: Demonstrate the importance of key performance areas in strategic management and design appropriate organizational structures and possess an ability to conceive various organizational dynamics.
    CO3: Compare and contrast early and contemporary theories of motivation and select and implement the right leadership practices in organizations that would enable systems orientation.
    CO4: Demonstrate an understanding on the usage and application of basic economic principles.
    CO5: Appreciate the various measures of macro-economic performance and interpret the prevailing economic health of the nation.
    """
    
    bloom_taxonomy_info = """
    Bloom's Taxonomy Levels:
    1. Remember: Recall facts and basic concepts
    2. Understand: Explain ideas or concepts
    3. Apply: Use information in new situations
    4. Analyze: Draw connections among ideas
    5. Evaluate: Justify a stand or decision
    6. Create: Produce new or original work
    """
    
    prompt = f"""
    Generate the following exam-style questions from this textbook content at {difficulty} difficulty level ({difficulty_desc}):

    1. 5 Multiple Choice Questions (with 4 options and correct answers)
    2. 5 Short Answer Questions (with 2–3 lines answers)
    3. 3 Long Answer Questions (for 6–8 marks, with detailed model answers)
    4. 4 Numerical Questions (using formulas, with step-by-step solutions)

    For each question:
    1. Assign the most relevant Course Outcome (CO) number from the list below:
    {course_outcomes}
    
    2. Specify the Bloom's Taxonomy level that best matches the cognitive skill required:
    {bloom_taxonomy_info}
    
    Content:
    {content}
    
    {"Return the result as a valid JSON object with the following structure:" if format_json else ""}
    
    {'''
    {
        "multipleChoice": [
            {
                "question": "Question text",
                "options": ["A. Option 1", "B. Option 2", "C. Option 3", "D. Option 4"],
                "correctAnswer": "A. Option 1",
                "explanation": "Why this is the correct answer",
                "courseOutcome": "CO1",
                "bloomsLevel": "Remember"
            }
        ],
        "shortAnswer": [
            {
                "question": "Question text",
                "answer": "Model answer (2-3 lines)",
                "courseOutcome": "CO2",
                "bloomsLevel": "Understand"
            }
        ],
        "longAnswer": [
            {
                "question": "Question text",
                "answer": "Detailed model answer (6-8 marks worth)",
                "markingScheme": "Bullet points of what would earn marks",
                "courseOutcome": "CO3",
                "bloomsLevel": "Analyze"
            }
        ],
        "numerical": [
            {
                "question": "Numerical problem statement",
                "solution": "Step-by-step solution",
                "answer": "Final numerical answer with units",
                "courseOutcome": "CO4",
                "bloomsLevel": "Apply"
            }
        ]
    }
    ''' if format_json else ""}
    
    Make sure each question is correctly tagged with both a course outcome (CO1-CO5) and a Bloom's taxonomy level.
    """
    
    try:
        response = model.generate_content(prompt)
        if format_json:
            response_text = response.text
            
            try:
                json_pattern = r'(\{[\s\S]*\})'
                json_matches = re.findall(json_pattern, response_text)
                
                if json_matches:
                    for potential_json in json_matches:
                        try:
                            return json.loads(potential_json)
                        except json.JSONDecodeError:
                            continue
                
                return json.loads(response_text)
            except json.JSONDecodeError:
                print("Warning: Could not parse response as JSON. Returning raw text.")
                return {"raw_response": response_text}
        else:
            return response.text
    except Exception as e:
        print(f"Error generating questions: {e}")
        return None

def save_to_json(data, output_path):
    """Save the generated questions to a JSON file."""
    try:
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(data, f, indent=4, ensure_ascii=False)
        print(f"Question bank saved to {output_path}")
        return True
    except Exception as e:
        print(f"Error saving JSON file: {e}")
        return False

def extract_text_from_file(file_path, page_range=None):
    """Extract text from various file formats."""
    file_ext = os.path.splitext(file_path)[1].lower()
    
    if file_ext == '.pdf':
        return extract_text_from_pdf(file_path, page_range)
    elif file_ext == '.docx':
        return extract_text_from_docx(file_path)
    elif file_ext in ['.pptx', '.ppt']:
        return extract_text_from_pptx(file_path)
    elif file_ext in ['.txt', '.text']:
        return extract_text_from_txt(file_path)
    else:
        print(f"Unsupported file format: {file_ext}")
        return None

def split_content_into_chunks(content, max_chunk_size=12000):
    """Split content into chunks of approximately max_chunk_size characters."""
    if len(content) <= max_chunk_size:
        return [content]
    
    paragraphs = re.split(r'\n\s*\n', content)
    
    paragraphs = [p for p in paragraphs if p.strip()]
    
    chunks = []
    current_chunk = []
    current_length = 0
    
    for paragraph in paragraphs:
        if current_length + len(paragraph) > max_chunk_size and current_chunk:
            chunks.append('\n\n'.join(current_chunk))
            current_chunk = []
            current_length = 0
        
        if len(paragraph) > max_chunk_size:
            sentences = re.split(r'(?<=[.!?])\s+', paragraph)
            
            for sentence in sentences:
                if current_length + len(sentence) > max_chunk_size and current_chunk:
                    chunks.append('\n\n'.join(current_chunk))
                    current_chunk = []
                    current_length = 0
                
                current_chunk.append(sentence)
                current_length += len(sentence) + 2 
                
        else:
            current_chunk.append(paragraph)
            current_length += len(paragraph) + 2 
    
    if current_chunk:
        chunks.append('\n\n'.join(current_chunk))
    
    return chunks

def process_chunk(args):
    """Process a single content chunk."""
    model, chunk, difficulty, format_json = args
    print(f"Processing chunk of {len(chunk)} characters...")
    return generate_question_bank(model, chunk, difficulty, format_json)

def process_chunks_parallel(model, chunks, difficulty, format_json, max_workers=2):
    """Process content chunks in parallel."""
    all_results = []
    
    with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
        futures = [executor.submit(process_chunk, (model, chunk, difficulty, format_json)) for chunk in chunks]
        
        for future in concurrent.futures.as_completed(futures):
            try:
                result = future.result()
                if result:
                    all_results.append(result)
            except Exception as e:
                print(f"Error processing chunk: {e}")
    
    if not all_results:
        return None
    
    if len(all_results) == 1:
        return all_results[0]
    
    merged = {
        "multipleChoice": [],
        "shortAnswer": [],
        "longAnswer": [],
        "numerical": []
    }
    
    for result in all_results:
        if isinstance(result, dict):
            if "raw_response" in result:
                print("Warning: One chunk returned raw text instead of JSON. Skipping.")
                continue
                
            for category in merged.keys():
                if category in result and isinstance(result[category], list):
                    merged[category].extend(result[category])
    
    max_counts = {
        "multipleChoice": 5,
        "shortAnswer": 5, 
        "longAnswer": 3,
        "numerical": 4
    }
    
    for category, max_count in max_counts.items():
        if len(merged[category]) > max_count:
            print(f"Limiting {category} questions to {max_count}")
            merged[category] = merged[category][:max_count]
    
    return merged

def process_chunks_sequentially(model, chunks, difficulty, format_json):
    """Process content chunks sequentially and merge results."""
    all_results = []
    
    for i, chunk in enumerate(chunks):
        print(f"Processing chunk {i+1}/{len(chunks)} ({len(chunk)} characters)...")
        result = generate_question_bank(model, chunk, difficulty, format_json)
        if result:
            all_results.append(result)
    
    if not all_results:
        return None
    
    if len(all_results) == 1:
        return all_results[0]
    
    merged = {
        "multipleChoice": [],
        "shortAnswer": [],
        "longAnswer": [],
        "numerical": []
    }
    
    for result in all_results:
        if isinstance(result, dict):
            if "raw_response" in result:
                print("Warning: One chunk returned raw text instead of JSON. Skipping.")
                continue
                
            for category in merged.keys():
                if category in result and isinstance(result[category], list):
                    merged[category].extend(result[category])
    
    max_counts = {
        "multipleChoice": 5,
        "shortAnswer": 5, 
        "longAnswer": 3,
        "numerical": 4
    }
    
    for category, max_count in max_counts.items():
        if len(merged[category]) > max_count:
            print(f"Limiting {category} questions to {max_count}")
            merged[category] = merged[category][:max_count]
    
    return merged

def get_question_bank_json():
    """Return the question bank as JSON data (not writing to file)."""
    return question_bank_data

def parse_page_range(page_range_str):
    """Parse a page range string like '1-5' into start and end page numbers."""
    try:
        if '-' in page_range_str:
            start, end = map(int, page_range_str.split('-'))
            return (start, end)
        else:
            page = int(page_range_str)
            return (page, page)
    except ValueError:
        print(f"Invalid page range format: {page_range_str}")
        return None

def main():
    parser = argparse.ArgumentParser(description='Generate exam questions from textbook files')
    parser.add_argument('--file', required=True, help='Path to the textbook file (PDF, DOCX, PPTX, or TXT)')
    parser.add_argument('--api_key', required=True, help='Gemini API key')
    parser.add_argument('--output', default='question_bank.json', help='Output JSON file path')
    parser.add_argument('--chapters', nargs='*', help='Specific chapters to process (by name)')
    parser.add_argument('--pages', help='Page range to process (start end)')
    parser.add_argument('--difficulty', default='medium', choices=['easy', 'medium', 'hard', 'mixed'], 
                        help='Difficulty level of questions')
    parser.add_argument('--text_only', action='store_true', help='Output as plain text instead of JSON')
    parser.add_argument('--no_file_output', action='store_true', help='Do not write output to file, just return JSON')
    parser.add_argument('--max_chunk_size', type=int, default=12000, help='Maximum characters per content chunk')
    parser.add_argument('--max_workers', type=int, default=1, help='Maximum number of parallel workers for processing chunks')
    
    args = parser.parse_args()
    
    page_range = None
    if args.pages:
        page_range = parse_page_range(args.pages)
        if not page_range:
            return
    
    if not os.path.exists(args.file):
        print(f"Error: File {args.file} not found")
        return
    
    if not check_and_install_dependencies():
        return
    
    content = extract_text_from_file(args.file, page_range)
    
    if not content:
        print("Failed to extract content from the file")
        return
    
    if args.chapters:
        content = extract_chapters(content, args.chapters)
    
    model = configure_api(args.api_key)
    if not model:
        return
    
    if len(content) > args.max_chunk_size:
        print(f"Content is too large ({len(content)} characters). Splitting into chunks...")
        chunks = split_content_into_chunks(content, args.max_chunk_size)
        print(f"Split content into {len(chunks)} chunks")
        
        if args.max_workers > 1:
            print(f"Processing chunks in parallel with {args.max_workers} workers...")
            question_bank = process_chunks_parallel(model, chunks, args.difficulty, not args.text_only, args.max_workers)
        else:
            print("Processing chunks sequentially...")
            question_bank = process_chunks_sequentially(model, chunks, args.difficulty, not args.text_only)
    else:
        question_bank = generate_question_bank(model, content, args.difficulty, not args.text_only)
    
    global question_bank_data
    question_bank_data = question_bank
    
    if question_bank and not args.no_file_output:
        if args.text_only:
            output_file = args.output
            if output_file.endswith('.json'):
                output_file = output_file.replace('.json', '.txt')
            
            try:
                with open(output_file, 'w', encoding='utf-8') as f:
                    f.write(question_bank)
                print(f"Question bank saved to {output_file}")
            except Exception as e:
                print(f"Error saving text file: {e}")
        else:
            save_to_json(question_bank, args.output)
    
    return question_bank

def get_question_bank():
    """Function that can be called from a server to get the generated question bank."""
    global question_bank_data
    if 'question_bank_data' in globals():
        return question_bank_data
    else:
        return {"error": "Question bank not generated yet"}

if __name__ == "__main__":
    main()